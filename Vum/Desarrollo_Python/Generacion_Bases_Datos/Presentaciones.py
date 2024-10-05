import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk  # Para usar la barra de progreso
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Image
from reportlab.lib.styles import getSampleStyleSheet
import requests
from io import BytesIO
import os

# Función para seleccionar archivo de Excel
def seleccionar_archivo():
    global ruta_archivo_seleccionado  # Declaramos que usaremos la variable global
    ruta_archivo_seleccionado = filedialog.askopenfilename(
        title="Selecciona un archivo Excel",
        filetypes=[("Archivos de Excel", "*.xlsx *.xls")]
    )
    if ruta_archivo_seleccionado:
        entry_ruta_archivo.delete(0, tk.END)  # Limpiamos el campo
        entry_ruta_archivo.insert(0, ruta_archivo_seleccionado)  # Insertamos la ruta seleccionada
        print(f"Ruta seleccionada: {ruta_archivo_seleccionado}")  # Imprimimos la ruta en consola

# Función para crear la presentación PPTX
def crear_presentacion_pptx(df, output_name, progress_bar, progress_label, root):
    prs = Presentation()
    total_rows = len(df)
    progress_step = 50 / total_rows  # Incremento para cada fila (50% de la barra se asigna al PPTX)

    for index, row in df.iterrows():
        # Descargar la imagen
        image_url = row['Actividad_Comercial[VISUAL]']
        image_stream = None
        if pd.notna(image_url) and image_url.startswith('http'):
            try:
                response = requests.get(image_url)
                response.raise_for_status()
                image_stream = BytesIO(response.content)
            except requests.RequestException as e:
                print(f"Error al descargar la imagen desde {image_url}: {e}")

        # Crear nueva diapositiva
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Usar el layout en blanco
        
        # Título de la diapositiva
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = str(row['Frecuencia[RAZON SOCIAL]'])  # Título de la diapositiva
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Agregar imagen a la diapositiva
        if image_stream:
            slide.shapes.add_picture(image_stream, Inches(0.5), Inches(2), width=Inches(4))

        # Agregar atributos en un cuadro de texto
        text_box = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(3))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for col in df.columns:
            if col != 'Frecuencia[NOMBRE PDV]' and col != 'Actividad_Comercial[fot_Tipo1]':
                p = text_frame.add_paragraph()
                p.text = str(row[col])
                p.font.size = Pt(16)

        # Actualizar progreso
        progress_bar.step(progress_step)
        percentage = progress_bar['value']
        progress_label.config(text=f"{percentage:.0f}%")
        root.update_idletasks()

    # Guardar la presentación PPTX en la ruta por defecto
    pptx_file = os.path.join(os.getcwd(), f"{output_name}.pptx")
    prs.save(pptx_file)
    print(f"Presentación PPTX guardada en {pptx_file}")
    return pptx_file

# Función para crear el PDF directamente desde los datos
def crear_pdf_directo_desde_datos(df, output_name, progress_bar, progress_label, root):
    pdf_file = os.path.join(os.getcwd(), f"{output_name}.pdf")
    doc = SimpleDocTemplate(pdf_file, pagesize=letter)
    elements = []
    styles = getSampleStyleSheet()

    elements.append(Paragraph(f"Reporte: {output_name}", styles['Title']))

    total_rows = len(df)
    progress_step = 50 / total_rows  # 50% restante para el PDF

    for index, row in df.iterrows():
        image_url = row['Actividad_Comercial[VISUAL]']
        image = None
        if pd.notna(image_url) and image_url.startswith('http'):
            try:
                response = requests.get(image_url)
                response.raise_for_status()
                image_stream = BytesIO(response.content)
                image = Image(image_stream, 2 * inch, 2 * inch)
            except requests.RequestException as e:
                print(f"Error al descargar la imagen desde {image_url}: {e}")

        # Tabla con datos de la fila
        data = [
            ['Razón Social', row['Frecuencia[RAZON SOCIAL]']],
            ['Nombre PDV', row['Frecuencia[NOMBRE PDV]']],
            ['Tipo Actividad', row['Actividad_Comercial[fot_Tipo1]']],
            ['Marca', row['Actividad_Comercial[fot_marca]']],
        ]

        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))

        elements.append(table)

        if image:
            elements.append(image)

        elements.append(Paragraph("<br/><br/>", styles['Normal']))

        # Actualizar barra de progreso
        progress_bar.step(progress_step)
        percentage = progress_bar['value']
        progress_label.config(text=f"{percentage:.0f}%")
        root.update_idletasks()

    doc.build(elements)
    print(f"PDF generado en {pdf_file}")
    return pdf_file

# Función para generar PPTX y PDF, y abrir el cuadro de diálogo para elegir dónde guardarlos
def crear_y_guardar_presentacion_y_pdf():
    ruta_excel = entry_ruta_archivo.get()
    nombre_presentacion = entry_nombre_presentacion.get()

    if not ruta_excel:
        messagebox.showwarning("Advertencia", "Debes seleccionar un archivo de Excel.")
        return
    if not nombre_presentacion:
        messagebox.showwarning("Advertencia", "Debes ingresar un nombre para la presentación.")
        return

    # Leer los datos del archivo Excel
    try:
        df = pd.read_excel(ruta_excel)
    except FileNotFoundError:
        messagebox.showerror("Error", "No se encuentra el archivo Excel.")
        return
    except Exception as e:
        messagebox.showerror("Error", f"Error al leer el archivo Excel: {e}")
        return

    # Resetear barra de progreso y porcentaje
    progress_bar['value'] = 0
    progress_label.config(text="0%")

    # Deshabilitar botón mientras se realiza el proceso
    btn_crear_presentacion.config(state=tk.DISABLED)

    # Crear la presentación PPTX
    pptx_file = crear_presentacion_pptx(df, nombre_presentacion, progress_bar, progress_label, ventana)

    # Crear el PDF directamente
    pdf_file = crear_pdf_directo_desde_datos(df, nombre_presentacion, progress_bar, progress_label, ventana)

    # Guardar el archivo PPTX en una ubicación específica
    ruta_guardado_pptx = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
    if ruta_guardado_pptx:
        os.rename(pptx_file, ruta_guardado_pptx)
        messagebox.showinfo("Guardado", f"Presentación PPTX guardada en {ruta_guardado_pptx}")

    # Guardar el archivo PDF en una ubicación específica
    ruta_guardado_pdf = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
    if ruta_guardado_pdf:
        os.rename(pdf_file, ruta_guardado_pdf)
        messagebox.showinfo("Guardado", f"PDF guardado en {ruta_guardado_pdf}")

    # Habilitar nuevamente el botón
    btn_crear_presentacion.config(state=tk.NORMAL)

# Función para crear y guardar la presentación PPTX
def crear_y_guardar_pptx():
    ruta_excel = entry_ruta_archivo.get()
    nombre_presentacion = entry_nombre_presentacion.get()

    if not ruta_excel:
        messagebox.showwarning("Advertencia", "Debes seleccionar un archivo de Excel.")
        return
    if not nombre_presentacion:
        messagebox.showwarning("Advertencia", "Debes ingresar un nombre para la presentación.")
        return

    # Leer los datos del archivo Excel
    try:
        df = pd.read_excel(ruta_excel)
    except FileNotFoundError:
        messagebox.showerror("Error", "No se encuentra el archivo Excel.")
        return
    except Exception as e:
        messagebox.showerror("Error", f"Error al leer el archivo Excel: {e}")
        return

    # Resetear barra de progreso y porcentaje
    progress_bar['value'] = 0
    progress_label.config(text="0%")

    # Generar la presentación PPTX
    pptx_file = crear_presentacion_pptx(df, nombre_presentacion, progress_bar, progress_label, ventana)

    # Guardar el archivo PPTX en una ubicación específica
    ruta_guardado_pptx = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
    if ruta_guardado_pptx:
        os.rename(pptx_file, ruta_guardado_pptx)
        messagebox.showinfo("Guardado", f"Presentación PPTX guardada en {ruta_guardado_pptx}")

# Función para crear y guardar el PDF
def crear_y_guardar_pdf():
    ruta_excel = entry_ruta_archivo.get()
    nombre_presentacion = entry_nombre_presentacion.get()

    if not ruta_excel:
        messagebox.showwarning("Advertencia", "Debes seleccionar un archivo de Excel.")
        return
    if not nombre_presentacion:
        messagebox.showwarning("Advertencia", "Debes ingresar un nombre para la presentación.")
        return

    # Leer los datos del archivo Excel
    try:
        df = pd.read_excel(ruta_excel)
    except FileNotFoundError:
        messagebox.showerror("Error", "No se encuentra el archivo Excel.")
        return
    except Exception as e:
        messagebox.showerror("Error", f"Error al leer el archivo Excel: {e}")
        return

    # Resetear barra de progreso y porcentaje
    progress_bar['value'] = 0
    progress_label.config(text="0%")

    # Generar el PDF
    pdf_file = crear_pdf_directo_desde_datos(df, nombre_presentacion, progress_bar, progress_label, ventana)

    # Guardar el archivo PDF en una ubicación específica
    ruta_guardado_pdf = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
    if ruta_guardado_pdf:
        os.rename(pdf_file, ruta_guardado_pdf)
        messagebox.showinfo("Guardado", f"PDF guardado en {ruta_guardado_pdf}")

# Función para cancelar el proceso
def cancelar_proceso():
    ventana.destroy()

# Crear la ventana principal
ventana = tk.Tk()
ventana.title("Creación de Presentación y PDF")
ventana.geometry("810x350")  # Ajustamos el tamaño de la ventana para incluir el pie de página

# Etiqueta y campo para seleccionar archivo de Excel
label_ruta_archivo = tk.Label(ventana, text="Ruta del archivo Excel:")
label_ruta_archivo.pack(pady=5)
entry_ruta_archivo = tk.Entry(ventana, width=50)
entry_ruta_archivo.pack(pady=5)
btn_seleccionar_archivo = tk.Button(ventana, text="Seleccionar Archivo", command=seleccionar_archivo)
btn_seleccionar_archivo.pack(pady=5)

# Etiqueta y campo para ingresar el nombre de la presentación
label_nombre_presentacion = tk.Label(ventana, text="Nombre de la Presentación:")
label_nombre_presentacion.pack(pady=5)
entry_nombre_presentacion = tk.Entry(ventana, width=50)
entry_nombre_presentacion.pack(pady=5)

# Barra de progreso
progress_bar = ttk.Progressbar(ventana, orient="horizontal", length=400, mode="determinate")
progress_bar.pack(pady=10)

# Etiqueta para mostrar el porcentaje
progress_label = tk.Label(ventana, text="0%", font=("Arial", 10))
progress_label.pack(pady=5)

# Botones para crear presentación y PDF, crear solo PPTX, o crear solo PDF
btn_crear_presentacion = tk.Button(ventana, text="Crear Presentación y PDF", command=crear_y_guardar_presentacion_y_pdf)
btn_crear_presentacion.pack(side=tk.LEFT, padx=20, pady=10)

btn_crear_pptx = tk.Button(ventana, text="Crear y Guardar PPTX", command=crear_y_guardar_pptx)
btn_crear_pptx.pack(side=tk.LEFT, padx=20, pady=10)

btn_crear_pdf = tk.Button(ventana, text="Crear y Guardar PDF", command=crear_y_guardar_pdf)
btn_crear_pdf.pack(side=tk.LEFT, padx=20, pady=10)

# Botón para cancelar el proceso
btn_cancelar = tk.Button(ventana, text="Cancelar Proceso", command=cancelar_proceso)
btn_cancelar.pack(side=tk.RIGHT, padx=50, pady=10)

# Pie de página con el texto "Desarrollado por el Equipo de Plataformas / Vision & Marketing"
footer_label = tk.Label(ventana, text="Desarrollado por el Equipo de Plataformas / Vision & Marketing", font=("Arial", 10), fg="gray")
footer_label.pack(side=tk.BOTTOM, fill=tk.X, pady=10)

# Ejecutar la ventana
ventana.mainloop()
